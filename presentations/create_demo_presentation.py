#!/usr/bin/env python3
"""
Create NEOlyzer demo presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor as RgbColor
import os

# Paths
SCRIPT_DIR = os.path.dirname(__file__)
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
CSS_LOGO_PATH = os.path.join(PROJECT_DIR, 'assets', 'CSS_logo_transparent.png')
OUTPUT_PATH = os.path.join(SCRIPT_DIR, 'NEOlyzer_Demo.pptx')

# Warning: regenerating will overwrite manual edits (animations, shadows, tweaks)
print("\n" + "="*70)
print("WARNING: Running this script will clobber the PowerPoint presentation")
print(f"         ({OUTPUT_PATH})")
print("         Any manual animations, shadows, or tweaks will be lost!")
print("="*70)
response = input("\nAre you sure you want to continue? [y/N]: ").strip().lower()
if response != 'y':
    print("Aborted.")
    exit(0)
print()

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors for professional theme
DARK_BLUE = RgbColor(0x1a, 0x23, 0x5b)  # Deep navy
MEDIUM_BLUE = RgbColor(0x2d, 0x3e, 0x83)  # Accent blue
LIGHT_BLUE = RgbColor(0x8a, 0xb4, 0xd9)  # Highlight (lighter shade)
WHITE = RgbColor(0xff, 0xff, 0xff)
GOLD = RgbColor(0xff, 0xc1, 0x07)  # Accent for highlights


def add_background(slide, color=DARK_BLUE):
    """Add solid background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, top, font_size=44, bold=True, color=WHITE):
    """Add a title text box"""
    left = Inches(0.5)
    width = Inches(12.333)
    height = Inches(1.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_body_text(slide, text, top, font_size=24, color=WHITE, left=0.75, width=11.833, center=False):
    """Add body text box"""
    txBox = slide.shapes.add_textbox(Inches(left), top, Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    if center:
        p.alignment = PP_ALIGN.CENTER
    return txBox


def add_bullet_points(slide, items, top, font_size=24, color=WHITE, left=0.75, line_spacing=1.0):
    """Add bulleted list with optional line spacing"""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(left), top, Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        if line_spacing != 1.0:
            p.line_spacing = line_spacing
    return txBox


def add_decorative_line(slide, top, color=GOLD):
    """Add a decorative accent line"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), top,
        Inches(5.333), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
slide_layout = prs.slide_layouts[6]  # Blank layout
title_slide = prs.slides.add_slide(slide_layout)
slide = title_slide  # Keep 'slide' reference for compatibility
add_background(slide)

# Main title
add_title_text(slide, "NEOlyzer", Inches(2), font_size=72, bold=True)

# Subtitle
add_title_text(slide, "A New Visualization Tool for Near-Earth Objects",
               Inches(3.2), font_size=36, bold=False, color=LIGHT_BLUE)

# Decorative line
add_decorative_line(slide, Inches(4.2))

# Authors
add_title_text(slide, "Rob Seaman", Inches(4.8), font_size=28, bold=True)
add_title_text(slide, "Catalina Sky Survey", Inches(5.3), font_size=22, bold=False, color=LIGHT_BLUE)

add_title_text(slide, "Claude (Anthropic)", Inches(5.9), font_size=28, bold=True)
add_title_text(slide, "AI Implementation Partner", Inches(6.4), font_size=22, bold=False, color=LIGHT_BLUE)


# =============================================================================
# SLIDE 2: The Prompt (Good morning...)
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide, MEDIUM_BLUE)

add_title_text(slide, "How This Presentation Was Made", Inches(0.4), font_size=40)
add_decorative_line(slide, Inches(1.1))

prompt_text = '''Good morning. It's been a few days since I worked on this. I need to create
a demo of NEOlyzer. Perhaps you can remind me - and yourself - of its high
points. Hmm. For that matter, I have both Powerpoint and Keynote on this
computer. Can you directly create a presentation while I kibitz?'''

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = prompt_text
p.font.size = Pt(26)
p.font.color.rgb = WHITE
p.font.italic = True

# Add context
add_body_text(slide,
    "This presentation was generated by Claude Code from a conversational prompt.\n\n"
    "NEOlyzer itself was developed through AI-assisted programming, achieving\n"
    "a Minimum Viable Product in a few days. Total effort was about two weeks.",
    Inches(4.8), font_size=22, color=LIGHT_BLUE)


# =============================================================================
# SLIDE 3: What is NEOlyzer?
# =============================================================================
slide3 = prs.slides.add_slide(slide_layout)
slide = slide3  # Keep compatibility
add_background(slide)

add_title_text(slide, "What is NEOlyzer?", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

bullets = [
    "Interactive visualization of the Near-Earth Object catalog",
    "Displays 40,000+ NEOs with smooth real-time animation",
    "Time range spanning 1550-2650 (using JPL DE440 ephemeris)",
    "Cross-platform: macOS, Linux, Windows (via WSL)",
    "Built with Python, PyQt6, matplotlib, Skyfield, SQLite"
]
add_bullet_points(slide, bullets, Inches(1.5), font_size=28, line_spacing=1.5)

# Add yellow asterisk after ephemeris line (4th bullet)
asterisk = slide.shapes.add_textbox(Inches(10.7), Inches(3.05), Inches(0.6), Inches(0.6))
tf = asterisk.text_frame
p = tf.paragraphs[0]
p.text = "*"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = GOLD

# Add comment box (rounded rectangle with yellow border, light blue fill)
# Centered horizontally: (13.333 - 7) / 2 = 3.17
comment_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.17), Inches(5.5),
    Inches(7), Inches(1.0)
)
comment_box.fill.solid()
comment_box.fill.fore_color.rgb = RgbColor(0xd0, 0xe8, 0xf8)  # Light blue background
comment_box.line.color.rgb = GOLD  # Yellow border
comment_box.line.width = Pt(3)

# Add text to comment box
tf = comment_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Orbits are not currently being integrated"
p.font.size = Pt(22)
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 4: Core Visualization
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Core Visualization Features", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

# Left column
left_title = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.6))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Map Projections"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

left_bullets = ["Rectangular (RA/Dec grid)", "Hammer (equal-area)",
                "Aitoff (azimuthal)", "Mollweide (pseudo-cylindrical)"]
add_bullet_points(slide, left_bullets, Inches(2.0), font_size=22, left=0.75)

# Right column
right_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.6))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "Coordinate Systems"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

right_bullets = ["Equatorial (standard obs)", "Ecliptic (orbital reference)",
                 "Galactic coordinates", "Opposition-centered"]
add_bullet_points(slide, right_bullets, Inches(2.0), font_size=22, left=7)

# Animation section (moved up)
anim_title = slide.shapes.add_textbox(Inches(0.75), Inches(4.0), Inches(11.5), Inches(0.6))
tf = anim_title.text_frame
p = tf.paragraphs[0]
p.text = "Real-Time Animation"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

anim_bullets = [
    "Variable playback rates (hours/days/months per second)",
    "Forward and backward playback",
    "~10 FPS with 40,000+ objects"
]
add_bullet_points(slide, anim_bullets, Inches(4.6), font_size=22, color=LIGHT_BLUE)

# Asterisk after "objects" line
asterisk = slide.shapes.add_textbox(Inches(4.6), Inches(5.3), Inches(0.5), Inches(0.5))
tf = asterisk.text_frame
p = tf.paragraphs[0]
p.text = "*"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = GOLD

# Comment box to the right
comment_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.5), Inches(6.3),
    Inches(6.0), Inches(0.8)
)
comment_box.fill.solid()
comment_box.fill.fore_color.rgb = RgbColor(0xd0, 0xe8, 0xf8)
comment_box.line.color.rgb = GOLD
comment_box.line.width = Pt(3)

tf = comment_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "As permitted by your hardware and OS"
p.font.size = Pt(20)
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 5: Data Encoding
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Multi-Dimensional Data Encoding", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

# Color encoding
color_title = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.6))
tf = color_title.text_frame
p = tf.paragraphs[0]
p.text = "Color Encodes:"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

color_bullets = ["V Magnitude (visual brightness)", "H Magnitude (intrinsic brightness)",
                 "CNEOS Discovery Site (observatory)"]
add_bullet_points(slide, color_bullets, Inches(2.0), font_size=22, left=0.75)

# Size encoding
size_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.6))
tf = size_title.text_frame
p = tf.paragraphs[0]
p.text = "Size Encodes:"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

size_bullets = ["V or H Magnitude", "Distance from Earth",
                "Earth MOID", "Orbital Period", "Eccentricity"]
add_bullet_points(slide, size_bullets, Inches(2.0), font_size=22, left=7)

# Filtering
filter_title = slide.shapes.add_textbox(Inches(0.75), Inches(4.3), Inches(11.5), Inches(0.6))
tf = filter_title.text_frame
p = tf.paragraphs[0]
p.text = "Filtering Options"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

filter_bullets = ["Dual magnitude limits (V and H, min and max)",
                  "NEO orbital classes (Atira, Aten, Apollo, Amor)",
                  "Earth MOID range filter",
                  "Hide objects before discovery date (lunation-based)"]
add_bullet_points(slide, filter_bullets, Inches(4.9), font_size=22)


# =============================================================================
# SLIDE 6: Analysis Tools
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Built-In Analysis Tools", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

# Heliocentric first (separate)
top_tool = ["Heliocentric Polar Chart - Sun-centered ecliptic view"]
add_bullet_points(slide, top_tool, Inches(1.5), font_size=26, line_spacing=1.5)

# Horizontal rule between sections
separator = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.75), Inches(2.3),
    Inches(11.5), Pt(2)
)
separator.fill.solid()
separator.fill.fore_color.rgb = GOLD
separator.line.fill.background()

# Rest of the tools
other_tools = [
    "MOID vs H Magnitude - Visualize Potentially Hazardous Asteroids",
    "NEO Discovery Timeline - Catalog growth through history",
    "Solar Elongation vs Distance - Observability planning",
    "Orbital Element Space (a vs e) - Classification view",
    "Lunar Phases Calendar - CLN tracking for observation scheduling"
]
add_bullet_points(slide, other_tools, Inches(2.5), font_size=26, line_spacing=1.5)

# Asterisk after MOID line (first line of other_tools, at y ~2.5)
asterisk = slide.shapes.add_textbox(Inches(11.5), Inches(2.5), Inches(0.5), Inches(0.5))
tf = asterisk.text_frame
p = tf.paragraphs[0]
p.text = "*"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = GOLD

# Comment box near bottom middle
comment_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.67), Inches(6.0),
    Inches(6), Inches(0.8)
)
comment_box.fill.solid()
comment_box.fill.fore_color.rgb = RgbColor(0xd0, 0xe8, 0xf8)
comment_box.line.color.rgb = GOLD
comment_box.line.width = Pt(3)

tf = comment_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Optimizing these charts TBD"
p.font.size = Pt(20)
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 7: Advanced Features
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Advanced Features", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

# Overlays
overlay_title = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(5.5), Inches(0.6))
tf = overlay_title.text_frame
p = tf.paragraphs[0]
p.text = "Astronomical Overlays"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

overlay_bullets = ["IAU constellation boundaries", "Bright star catalog",
                   "Ecliptic and Galactic planes", "Observer horizon & twilight"]
add_bullet_points(slide, overlay_bullets, Inches(2.0), font_size=22, left=0.75)

# Catalog comparison
catalog_title = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.6))
tf = catalog_title.text_frame
p = tf.paragraphs[0]
p.text = "Catalog Comparison"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

catalog_bullets = ["Load alternate catalog versions", "Blink between catalogs",
                   "View new/deleted/changed objects", "Pre-computed position cache"]
add_bullet_points(slide, catalog_bullets, Inches(2.0), font_size=22, left=7)

# Interactive features
interact_title = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.5), Inches(0.6))
tf = interact_title.text_frame
p = tf.paragraphs[0]
p.text = "Interactive Features"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GOLD

interact_bullets = [
    "Click any NEO for detailed orbital elements and discovery info",
    "Shift+Click for constellation ID",
    "Sortable data tables with CSV export"
]
add_bullet_points(slide, interact_bullets, Inches(5.1), font_size=22, color=LIGHT_BLUE)


# =============================================================================
# SLIDE 8: Performance & Architecture
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Performance Architecture", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

# HDF5 section with tight spacing for sub-items
hdf5_bullets = [
    "HDF5 position cache with variable precision tiers",
    "   - High precision (daily) for current year",
    "   - Medium precision (weekly) for ±5 years",
    "   - Low precision (monthly) for extended range",
]
add_bullet_points(slide, hdf5_bullets, Inches(1.5), font_size=24, line_spacing=1.0)

# Rest of performance bullets with wider spacing
other_perf_bullets = [
    "~10 FPS sustained for 40,000+ objects",
    "Magnitude hysteresis reduces visual twinkling",
    "Efficient LineCollection for boundary rendering",
    "SQLite database via SQLAlchemy ORM"
]
add_bullet_points(slide, other_perf_bullets, Inches(3.4), font_size=24, line_spacing=1.5)


# =============================================================================
# SLIDE 9: Future Directions
# =============================================================================
slide = prs.slides.add_slide(slide_layout)
add_background(slide)

add_title_text(slide, "Future Directions", Inches(0.3), font_size=44)
add_decorative_line(slide, Inches(1.0))

future_bullets = [
    "NEOfixer integration",
    "   - Additional observational planning modes",
    "Performance optimization for 100,000+ NEOs (catalog growing rapidly)",
    "Database schema normalization as features expand",
    "Community feedback and feature requests welcome"
]
add_bullet_points(slide, future_bullets, Inches(1.5), font_size=26, line_spacing=1.5)

# Contact
contact_title = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.5), Inches(0.6))
tf = contact_title.text_frame
p = tf.paragraphs[0]
p.text = "Contact: rseaman@arizona.edu"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GOLD


# =============================================================================
# SLIDE 10: Live Demo (last slide)
# =============================================================================
last_slide = prs.slides.add_slide(slide_layout)
slide = last_slide  # Keep 'slide' reference for compatibility
add_background(slide, MEDIUM_BLUE)

add_title_text(slide, "Live Demo", Inches(2.5), font_size=60)
add_decorative_line(slide, Inches(3.4))

demo_points = "Projections  |  Filtering  |  Catalog Blinking  |  Animation"
add_body_text(slide, demo_points, Inches(4.2), font_size=28, color=LIGHT_BLUE, center=True)


# =============================================================================
# Add CSS Logo to Title and Last Slides
# =============================================================================
def add_logo_with_circle(slide, logo_path, left, top, logo_size):
    """Add logo centered on a white circle background"""
    circle_size = logo_size * 1.15  # Circle slightly larger than logo
    circle_offset = (circle_size - logo_size) / 2

    # Add white circle first (so it's behind the logo)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(circle_size), Inches(circle_size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()  # No border

    # Add logo centered on circle
    slide.shapes.add_picture(
        logo_path,
        Inches(left + circle_offset), Inches(top + circle_offset),
        height=Inches(logo_size), width=Inches(logo_size)
    )

if os.path.exists(CSS_LOGO_PATH):
    # Title slide: larger logo, centered between left edge and "Claude" name
    # Claude name is centered text, so visually around x=6.67 (slide center)
    # Center logo between 0 and 6.67 = 3.33
    title_logo_size = 2.2  # Larger on title slide
    title_circle_size = title_logo_size * 1.15
    title_center_x = 2.33
    title_left = title_center_x - (title_circle_size / 2)
    title_top = 7.5 - 0.5 - title_circle_size  # 0.5 margin from bottom
    add_logo_with_circle(title_slide, CSS_LOGO_PATH, title_left, title_top, title_logo_size)

    # Last slide: vertically aligned with email, equal margin from right and bottom
    last_logo_size = 1.8
    last_circle_size = last_logo_size * 1.15
    email_middle_y = 5.2 + 0.3  # Middle of email text box
    last_top = email_middle_y - (last_circle_size / 2)
    # Bottom margin = slide_height - (top + circle_size)
    bottom_margin = 7.5 - (last_top + last_circle_size)
    # Set right margin equal to bottom margin
    last_left = 13.333 - bottom_margin - last_circle_size
    add_logo_with_circle(last_slide, CSS_LOGO_PATH, last_left, last_top, last_logo_size)

    print(f"Added CSS logo with white circle from: {CSS_LOGO_PATH}")
else:
    print(f"Warning: CSS logo not found at {CSS_LOGO_PATH}")


# =============================================================================
# Save presentation
# =============================================================================
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
